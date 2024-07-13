from docx import Document
import pptx
import re
import pdfplumber


class naive_text_extraction:
    def __init__(self):
        self.file_instance = None
        self.extracted_text = {}
        self.line_num_threshold = 2
        self.space_num_threshold = 2


    def start(self, doc):        
        self.extracted_text = {}
        self.extracted_text["document_type"] = doc.name.split(".")[-1].lower()
        
        if self.extracted_text["document_type"] not in ["pdf", "docx", "pptx", "ppt"]:
            return None
        
        self.file_instance = doc
        self.extracted_text["contents"] = {}
        self.extracted_text["document_title"] = doc.name

        if self.extracted_text["document_type"] == "pdf":
            self.parse_pdf()
        elif self.extracted_text["document_type"] == "docx":
            self.parse_docx()
        elif self.extracted_text["document_type"] in ["pptx", "ppt"]:
            self.parse_ppt()   

        extracted_text = self.extracted_text
        self.extracted_text = {}
        self.file_instance = None

        return extracted_text
        

    def parse_pdf(self):    
        self.extracted_text["contents_index_type"] = "page_number"
        with pdfplumber.open(self.file_instance) as pdf:
            for page in pdf.pages:
                page_num = page.page_number
                text_blocks = page.extract_text()
                if text_blocks:
                    page_text = self.remove_excessive_lines_and_spaces(text_blocks)
                else:
                    page_text = "" 
                self.extracted_text["contents"][page_num] = {"text": page_text}


    def parse_docx(self):
        doc = Document(self.file_instance)

        self.extracted_text["contents_index_type"] = "whole_text_or_table_number"

        cnt_text_idx = 1
        cnt_table_idx = 1
        
        full_text = ""
        for para in doc.paragraphs:
            text = para.text
            if text in ["", " ", None]:
                continue
            full_text += text + "\n"

        full_text = self.remove_excessive_lines_and_spaces(full_text)
        self.extracted_text["contents"][cnt_text_idx] = {"text": full_text}

        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(",".join(row_data))
            csv_table = "\n".join(table_data)
            
            if cnt_table_idx in self.extracted_text["contents"]:
                self.extracted_text["contents"][cnt_table_idx]["table"] = csv_table
            else:
                self.extracted_text["contents"][cnt_table_idx] = {"table": csv_table}
            cnt_table_idx += 1


    def parse_ppt(self):
        presentation = pptx.Presentation(self.file_instance)

        self.extracted_text["contents_index_type"] = "slide_number"
        
        for slide in presentation.slides:
            slide_number = presentation.slides.index(slide) + 1
            self.extracted_text["contents"][slide_number] = {}
            texts = ""
            tables = ""
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text = ""
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + "\n"
                    texts += text + "\n\n"
                
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text_frame.text if cell.text_frame else ""
                            row_data.append(cell_text)
                        table_data.append(",".join(row_data))
                    csv_table = "\n".join(table_data)
                    tables += csv_table + "\n\n"
            if texts:
                texts = self.remove_excessive_lines_and_spaces(texts)
                self.extracted_text["contents"][slide_number]["text"] = texts
            if tables:
                self.extracted_text["contents"][slide_number]["table"] = tables


    def remove_excessive_lines_and_spaces(self, text):
        # Replace 3 or more consecutive newlines with 3 newlines
        text = re.sub(r'\n{'+ str(self.line_num_threshold) + ',}', '\n'*self.line_num_threshold, text)
        # Replace 3 or more consecutive spaces with 3 spaces
        text = re.sub(r' {'+ str(self.space_num_threshold) + ',}', ' '*self.space_num_threshold, text)

        return text

















